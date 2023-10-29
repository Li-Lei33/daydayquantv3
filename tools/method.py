import datetime as dt
import json
import pandas as pd
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches

import seaborn as sns
from pptx import Presentation
from pptx.util import Inches

import os
import numpy as np
import csv
from io import BytesIO

from config.constant import Model

plt.rcParams['font.sans-serif'] = ['SimHei']  # 例如使用黑体
plt.rcParams['axes.unicode_minus'] = False  # 解决保存图像时负号'-'显示为方块的问题


def bars2ls(bars, interval, start, end):
    # bar: [datetime, open, high, low, close, volume]
    start_dt = dt.datetime.strptime(start, "%Y-%m-%d")
    end_dt = dt.datetime.strptime(end, "%Y-%m-%d")

    bar_ls = []
    started = False
    is_first = False
    is_last = False
    for bar in bars:
        bar[0] = dt.datetime.strptime(bar[0], '%Y-%m-%d %H:%M:%S')
        if bar[0] < start_dt or bar[0] > end_dt:
            continue
        if bar[0].minute % interval == 0:
            started = True
            is_first = True
            is_last = False
        elif (bar[0].minute+1) % interval == 0:
            is_first = False
            is_last = True
        else:
            is_first = False
            is_last = False
            pass
        if not started:
            continue
        if is_first:
            datetime_val = bar[0]
            open_val = bar[1]
            high_val = bar[2]
            low_val = bar[3]
            close_val = bar[4]
            volume_val = bar[5]
        elif not is_first and not is_last:
            datetime_val = datetime_val
            open_val = open_val
            high_val = bar[2] if bar[2] > high_val else high_val
            low_val = bar[3] if bar[3] < low_val else low_val
            close_val = bar[4]
            volume_val += bar[5]
        elif is_last:
            datetime_val = datetime_val
            open_val = open_val
            high_val = bar[2] if bar[2] > high_val else high_val
            low_val = bar[3] if bar[3] < low_val else low_val
            close_val = bar[4]
            volume_val += bar[5]
            # bar_ls.append([datetime_val, open_val, high_val, low_val, close_val, volume_val, name_val])
            # bar_ls.append([open_val, high_val, low_val, close_val, volume_val])
            # bar_ls.append([open_val, high_val, low_val, close_val])
            bar_ls.append([datetime_val, open_val, high_val, low_val, close_val, volume_val])
        else:
            pass
    return bar_ls


def list_model_directories(project_path):
    """
    List directories in the specified path that start with 'model'

    :param path: str, path to the directory to search in
    :return: list of directory names that start with 'model'
    """
    # 确保提供的路径是一个存在的目录
    if not os.path.isdir(project_path):
        print(f"'{project_path}' is not a directory or does not exist.")
        return []

    # 获取path目录下的所有文件和文件夹
    all_items = os.listdir(project_path)

    # 过滤出以'model'开头且是目录的项
    model_dirs = [item for item in all_items
                  if item.startswith('model') and os.path.isdir(os.path.join(project_path, item))]

    # 如果conf配置里的now_model已经被删了，更新为现有model的最后一个
    conf = read_json(os.path.join(project_path, "config", "conf.json"))
    now_model_name = Model.MODEL_PRE.value + str(conf["now_model"])
    if now_model_name not in model_dirs:
        conf["now_model"] = int(model_dirs[-1].split(Model.MODEL_PRE.value)[-1])
        save_json(conf, os.path.join(project_path, "config", "conf.json"))

    return model_dirs


def read_json(path):
    """输入路径读取json文件，返回对应dict"""
    import json
    with open(path, 'r') as file:
        dc = json.load(file)
    return dc


def save_json(dc, path):
    import json
    with open(path, 'w') as file:
        json.dump(dc, file)
    return dc


def dt_str():
    dt_str = dt.datetime.now().strftime("[%H:%M:%S]: ")
    return dt_str


def gen_ppt_t1(result_dir_path):
    """按1方案根据回测结果生成回测报告"""
    # 回测报告ppt生成代码，可封装为单独的函数，供给出不同风格的回测报告

    # 获取result中的结果数据，生成一份ppt，作为回测报告
    # 获取文件夹下所有result文件
    files_and_folders = os.listdir(result_dir_path)
    # 过滤出文件（仅需要回测产生的csv文件，其他文件过滤，否则报错）
    files = [f for f in files_and_folders if os.path.isfile(os.path.join(result_dir_path, f)) and f.endswith(".csv")]

    # 预创建PPT
    prs = Presentation()
    # 遍历文件
    for file_one in files:
        file_path = os.path.join(result_dir_path, file_one)
        # 读取文件内容
        result_ls = []
        with open(file_path, mode='r', newline='') as file:
            csv_reader = csv.reader(file)
            for row in csv_reader:
                result_ls.append(row)
        # 解析文件内容为两部分：参数对，净值和交易行为
        params_dc = {}
        cash_value_ls = []
        order_record_ls = []
        trade_record_ls = []
        for one in result_ls:
            mark = one[0]
            if mark == '0':
                params_dc[one[1]] = one[2]
            elif mark == '1':
                # [cash ,value, close, datetime]
                cash_value_ls.append([float(one[1]), float(one[2]), float(one[3]), one[4]])
            elif mark == '2':
                # [对应K线索引, 订单价格, +-订单数量]
                order_record_ls.append([len(cash_value_ls), float(one[1]), float(one[2])])
            elif mark == '3':
                # [对应K线索引, 交易状态, 收益率]
                trade_record_ls.append([len(cash_value_ls), one[1], float(one[2])])
            else:
                raise ValueError('未知的标记')
        # 回测分析报告
        # （弃用）收益率、无成本收益率、交易次数（次/日）、夏普指数（变式）、图：现金-净值、收盘价、最大回撤
        # 1、收益率（净） 2、交易成本占比 3、最大回撤 4、上涨时间：下跌时间 5、盈利单：亏损单
        # 收益率
        end_time = dt.datetime.strptime(cash_value_ls[-1][3], "%Y-%m-%d %H:%M")
        start_time = dt.datetime.strptime(cash_value_ls[0][3], "%Y-%m-%d %H:%M")
        delta = end_time - start_time
        years_diff = delta.days / 365.25 if delta.days != 0 else 1 / 365.25
        power = 1 / years_diff
        earning_rate = (cash_value_ls[-1][1] - cash_value_ls[0][1]) / cash_value_ls[0][1]
        # 避免基数为负数
        earning_rate_year = pow(earning_rate, power) if earning_rate >= 0 else -pow(-earning_rate, power)
        earning_rate_year_str = f"{earning_rate_year:.2%}"
        # 交易成本占比
        trade_vol = 0
        for o in order_record_ls:
            trade_vol += o[1] * abs(o[2])
        trade_cost = trade_vol * 0.002
        trade_cost_rate = trade_cost / cash_value_ls[0][1]
        trade_cost_rate_str = f"{trade_cost_rate:.2%}"
        # 最大回撤，最大回撤图
        max_withdraw_ls = []
        max_value = 0
        for o in cash_value_ls:
            now_value = o[2]
            max_value = now_value if now_value > max_value else max_value
            max_withdraw_ls.append((max_value - now_value) / max_value)
        max_withdraw = max(max_withdraw_ls)
        max_withdraw_str = f"{max_withdraw:.2%}"
        # 上涨下跌时间
        up_num = 0
        down_num = 0
        for i in range(len(cash_value_ls) - 1):
            if cash_value_ls[i + 1][2] - cash_value_ls[i][2] > 0:
                up_num += 1
            else:
                down_num += 1
        up_down_rate = up_num / down_num if down_num != 0 else 1e8
        up_down_rate_str = f"{up_down_rate:.2f}"
        # 盈利亏损单
        profit_num = 0
        loss_num = 0
        for i in range(len(trade_record_ls)):
            if trade_record_ls[i][1] == 'Close':
                if trade_record_ls[i][2] >= 0:
                    profit_num += 1
                else:
                    loss_num += 1
        profit_loss_rate = profit_num / loss_num if loss_num != 0 else 1e8
        profit_loss_rate_str = f"{profit_loss_rate:.2f}"

        # 生成一页ppt
        # 绘制cash-value图、最大回撤图、close图
        x = list(range(len(cash_value_ls)))
        cash_ls = []
        value_ls = []
        close_ls = []
        for o in cash_value_ls:
            cash_ls.append(o[0])
            value_ls.append(o[1])
            close_ls.append(o[2])

        # 创建画布，一行三列
        fig, axes = plt.subplots(nrows=3, ncols=1, figsize=(8, 15))
        # 绘制value图
        # sns.lineplot(x=x, y=cash_ls, ax=axes[0])
        sns.lineplot(x=x, y=value_ls, ax=axes[0])
        axes[0].set_title('value')

        sns.lineplot(x=x, y=max_withdraw_ls, ax=axes[1])
        axes[1].set_title('max withdraw')

        sns.lineplot(x=x, y=close_ls, ax=axes[2])
        axes[2].set_title('close')

        plt.tight_layout()
        # plt.show()  # ---------------- 1
        # 保存Matplotlib图表为图片
        img_stream_1 = BytesIO()
        plt.savefig(img_stream_1, format="png")
        img_stream_1.seek(0)

        # 该数据最相关的图为五维图，但在表格图中需要，故提前设置
        label_ls = ["收益率（净）", "交易成本占比", "最大回撤", "涨跌比", "盈亏比"]
        full_score_val = np.array([2, 0, 0.05, 5, 5])
        zero_score_val = np.array([-1, 0.3, 1, 0, 0])
        score_attr = np.array([1, -1, -1, 1, 1])  # 用来标注该参数的属性，越大越好or越小越好

        # 重新创建变量，美化代码
        full = full_score_val
        zero = zero_score_val

        # 绘制指标列表
        data_tab = [
            ['指标名', '数值', '区间（零分 ： 满分）', '备注'],
            ['收益率（静）', earning_rate_year_str, f'{zero[0]} : {full[0]}', '考虑交易成本，折合年化'],
            ['交易成本占比', trade_cost_rate_str, f'{zero[1]} : {full[1]}', ''],
            ['最大回撤', max_withdraw_str, f'{zero[2]} : {full[2]}', ''],
            ['涨跌比', up_down_rate_str, f'{zero[3]} : {full[3]}', '涨跌时间比，排除0变动时间'],
            ['盈亏比', profit_loss_rate_str, f'{zero[4]} : {full[4]}', '盈亏交易数量比']
        ]
        fig, ax = plt.subplots()

        table = ax.table(cellText=data_tab, loc='center', colWidths=[0.2, 0.15, 0.3, 0.4])
        table.scale(1, 2)  # 调整行高
        # 隐藏坐标轴
        ax.axis('off')
        # 设置表格样式
        table.auto_set_font_size(False)
        table.set_fontsize(11)
        table.scale(1.2, 1.2)  # 调整表格大小

        # plt.show()  # ---------------- 2
        # 保存Matplotlib图表为图片
        img_stream_2 = BytesIO()
        plt.savefig(img_stream_2, format="png")
        img_stream_2.seek(0)

        # 绘制五维图
        # 1、收益率（净）
        # 2、交易成本占比
        # 3、最大回撤
        # 4、上涨时间：下跌时间
        # 5、盈利单：亏损单

        def press(real, full_val, zero_val, attr):
            """返回截断后数据，优于full_score等于full_score，劣于zero_score等于zero_score，其他不变"""
            if real * attr > full_val * attr:
                return full_val
            elif real * attr < zero_val * attr:
                return zero_val
            else:
                return real

        earning_rate_year_press = press(earning_rate_year, full_score_val[0], zero_score_val[0], score_attr[0])
        trade_cost_rate_press = press(trade_cost_rate, full_score_val[1], zero_score_val[1], score_attr[1])
        max_withdraw_press = press(max_withdraw, full_score_val[2], zero_score_val[2], score_attr[2])
        up_down_rate_press = press(up_down_rate, full_score_val[3], zero_score_val[3], score_attr[3])
        profit_loss_rate_press = press(profit_loss_rate, full_score_val[4], zero_score_val[4], score_attr[4])

        real_value = np.array([earning_rate_year_press, trade_cost_rate_press, max_withdraw_press, up_down_rate_press, profit_loss_rate_press])
        # 将real_value转为0-1之间的数
        # real_value到zero_score之间的值
        distance_real_zero = (real_value - zero_score_val) * score_attr
        # full_score到zero_score之间的值
        distance_full_zero = (full_score_val - zero_score_val) * score_attr
        index_value = distance_real_zero / distance_full_zero

        # 定义五边形的顶点数
        num_vertices = 5

        # 计算五边形顶点的位置
        angles_ls = np.linspace(0, 2 * np.pi, num_vertices, endpoint=False).tolist()

        # 完成雷达图的封闭
        stats = np.concatenate((index_value, [index_value[0]]))
        angles = np.concatenate((angles_ls, [angles_ls[0]]))

        # 创建画布
        fig, ax = plt.subplots(figsize=(6, 6), subplot_kw=dict(polar=True))

        # # 绘制五边形轮廓
        # ax.plot(x, y, 'b-')

        # 填充五边形颜色
        ax.fill(angles, stats, color='red', alpha=0.25)

        # 在雷达图上添加标签
        ax.set_yticklabels([])
        ax.set_xticks(angles[:-1])
        ax.set_xticklabels(label_ls)

        # 添加图例
        red_patch = mpatches.Patch(color='red', alpha=0.25, label='纬度值')
        plt.legend(handles=[red_patch])

        # 保存Matplotlib图表为图片
        img_stream_3 = BytesIO()
        plt.savefig(img_stream_3, format="png")
        img_stream_3.seek(0)

        # 添加一张幻灯片
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # 第一张幻灯片加上bt回测结果图像
        fig_path = os.path.join(result_dir_path, "result figure.png")
        # 图片位置  -------------- 0
        left = Inches(0.5)  # 调整图片位置
        top = Inches(0.5)
        width = Inches(9)  # 图片宽度
        height = Inches(6)  # 图片高度
        # 在幻灯片中插入图片
        slide.shapes.add_picture(fig_path, left, top, width=width, height=height)

        # 添加一张幻灯片
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 图片位置  -------------- 1
        left = Inches(0.5)  # 调整图片位置
        top = Inches(0.5)
        width = Inches(9)  # 图片宽度
        height = Inches(6)  # 图片高度
        # 在幻灯片中插入图片
        slide.shapes.add_picture(img_stream_1, left, top, width=width, height=height)

        # 添加一张幻灯片
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 图片位置  -------------- 2
        left = Inches(0.5)  # 调整图片位置
        top = Inches(1)
        width = Inches(4)  # 图片宽度
        height = Inches(4)  # 图片高度
        # 在幻灯片中插入图片
        slide.shapes.add_picture(img_stream_2, left, top, width=width, height=height)

        # # 添加一张幻灯片
        # slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 图片位置  -------------- 3
        left = Inches(5)  # 调整图片位置
        top = Inches(1)
        width = Inches(4)  # 图片宽度
        height = Inches(4)  # 图片高度
        # 在幻灯片中插入图片
        slide.shapes.add_picture(img_stream_3, left, top, width=width, height=height)

    # 循环结束，保存幻灯片
    pptx_name = "presentation.pptx"
    prs.save(os.path.join(result_dir_path, pptx_name))
    print(f"{dt_str()}回测报告路径：{os.path.join(result_dir_path, pptx_name)}")
    # windows下
    os.system(f"start {os.path.join(result_dir_path, pptx_name)}")
    # mac下
    # os.system(f"open {os.path.join(result_dir_path, pptx_name)}")

    pass


def gen_ppt_f1(data_factor: pd.DataFrame, result_dir_path):
    """接收行情数据和因子数据的DataFrame，并给出分析报告"""
    # 因子相关性热力图
    # 提取因子名，作为列表
    factor_ls = []
    for f in data_factor.columns.values:
        # 不属于列表中的字段名均视为因子，请规范命名
        if f not in ['date', 'datetime', 'money', 'open_interest', 'contract', 'open', 'high', 'low', 'close', 'volume']:
            factor_ls.append(f)
    # 提取因子值
    factor_value_df = data_factor[factor_ls]
    # 创建ppt
    prs = Presentation()
    n = 0
    draw_factor_ls = []
    for f in factor_ls:
        print(f"{dt_str()}开始计算因子：{f}")
        n += 1
        # 给出因子特征，平均值、中位数、最大值、最小值、并将其附到因子正态分布图下方
        # 计算平均值
        mean_value = factor_value_df[f].mean()
        # 计算中位数
        median_value = factor_value_df[f].median()
        # 计算最大值
        max_value = factor_value_df[f].max()
        # 计算最小值
        min_value = factor_value_df[f].min()
        # 设置图形的大小
        plt.figure(figsize=(10, 6))
        # 过滤数据中超出上下限的值
        # 逐个添加过滤规则
        pd.options.mode.chained_assignment = None  # 默认为 'warn'，确定操作无误，忽略警告信息
        if f == 'mom':
            factor_value_df.loc[:, f] = factor_value_df[f].clip(lower=-300, upper=300).copy()
        elif f == 'roc':
            factor_value_df.loc[:, f] = factor_value_df[f].clip(lower=-2, upper=2).copy()
        else:
            pass

        # 绘制直方图和KDE
        ax = sns.histplot(factor_value_df[f], kde=True, bins=30, kde_kws={})
        # 在每个柱体上添加频数
        for p in ax.patches:
            ax.annotate(f'{int(p.get_height())}',
                        (p.get_x() + p.get_width() / 2., p.get_height()),
                        ha='center', va='center',
                        xytext=(0, 9),
                        textcoords='offset points')

        # 添加标题和标签
        plt.title(f'因子({f}) 特征: -mean({round(mean_value, 3)}) -median({round(median_value, 3)}) -max({round(max_value, 3)}) -min({round(min_value, 3)})', fontsize=16)
        plt.xlabel('因子值')
        plt.ylabel('频数')

        # 显示图形
        # plt.show()
        img_stream = BytesIO()
        plt.savefig(img_stream, format='png')
        img_stream.seek(0)
        # 添加一张幻灯片
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        left = Inches(0.5)  # 调整图片位置
        top = Inches(0.5)
        width = Inches(9)  # 图片宽度
        height = Inches(6)  # 图片高度
        # 在幻灯片中插入图片
        slide.shapes.add_picture(img_stream, left, top, width=width, height=height)

        # 展示局部（随机1000个节点）因子（五个，超过五个另起一图）及对应收盘价图
        if n % 3 == 0:
            draw_factor_ls.append(f)

            draw_factor_ls.append('close')
            # 获取绘图数据
            draw_data = data_factor[draw_factor_ls]
            # 绘图
            # 创建一个子图网格，其中每个子图共用x轴
            fig, axes = plt.subplots(nrows=len(draw_data.columns), sharex=True, figsize=(16, 12))

            # 对于DataFrame的每一列，绘制一个子图
            for i, col in enumerate(draw_data.columns):
                # draw_data[col].plot(ax=axes[i], title=col)
                axes[i].scatter(draw_data[col].index, draw_data[col].values, s=1)
                axes[i].set_title(col, fontsize=20)

                axes[i].set_ylabel('Value', fontsize=20)

            # 设置x轴标签
            axes[-1].set_xlabel('Index', fontsize=20)

            # 调整子图之间的间距
            plt.tight_layout()

            # 显示图形
            # plt.show()
            img_stream = BytesIO()
            plt.savefig(img_stream, format='png')
            img_stream.seek(0)
            # 添加一张幻灯片
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            left = Inches(0.5)  # 调整图片位置
            top = Inches(0.5)
            width = Inches(9)  # 图片宽度
            height = Inches(6)  # 图片高度
            # 在幻灯片中插入图片
            slide.shapes.add_picture(img_stream, left, top, width=width, height=height)

            # 重置绘图因子列表
            draw_factor_ls = []
        else:
            draw_factor_ls.append(f)
    # for循环结束，若仍存在需要绘制的因子
    if len(draw_factor_ls) != 0:
        draw_factor_ls.append('close')
        # 获取绘图数据
        draw_data = data_factor[draw_factor_ls]
        # 绘图
        # 创建一个子图网格，其中每个子图共用x轴
        fig, axes = plt.subplots(nrows=len(draw_data.columns), sharex=True, figsize=(8, 6))

        # 对于DataFrame的每一列，绘制一个子图
        for i, col in enumerate(draw_data.columns):
            draw_data[col].plot(ax=axes[i], title=col)
            axes[i].set_ylabel('Value')

        # 设置x轴标签
        axes[-1].set_xlabel('Index')

        # 调整子图之间的间距
        plt.tight_layout()

        # 显示图形
        # plt.show()
        img_stream = BytesIO()
        plt.savefig(img_stream, format='png')
        img_stream.seek(0)
        # 添加一张幻灯片
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        left = Inches(0.5)  # 调整图片位置
        top = Inches(0.5)
        width = Inches(9)  # 图片宽度
        height = Inches(6)  # 图片高度
        # 在幻灯片中插入图片
        slide.shapes.add_picture(img_stream, left, top, width=width, height=height)

        # 重置绘图因子列表
        draw_factor_ls = []
    else:
        pass

    # 循环结束，保存幻灯片
    pptx_name = "factor-feature.pptx"
    prs.save(os.path.join(result_dir_path, pptx_name))
    print(f"{dt_str()}pptx文件已生成，路径为{os.path.join(result_dir_path, pptx_name)}。\n"
          f"如需调整绘图数据，请至gen_ppt_f1方法中添加过滤规则。")

    # windows下
    os.system(f"start {os.path.join(result_dir_path, pptx_name)}")
    # mac下
    # os.system(f"open {os.path.join(result_dir_path, pptx_name)}")


